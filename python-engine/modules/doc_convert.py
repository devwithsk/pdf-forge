import os
import sys
import json
import shutil
import subprocess
import tempfile
from pdf2docx import Converter
from reportlab.lib.pagesizes import letter, landscape
from limits import assert_pdf_limits
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

# Try to import Windows COM libraries for high-fidelity conversion if running on Windows with Office installed
COM_AVAILABLE = False
if sys.platform == "win32":
    try:
        import win32com.client
        import comtypes.client
        COM_AVAILABLE = True
    except ImportError:
        pass

# Probe for headless LibreOffice once at startup — avoids repeated shell calls per request
LIBREOFFICE_AVAILABLE = shutil.which('libreoffice') is not None or shutil.which('soffice') is not None
_LO_BIN = shutil.which('libreoffice') or shutil.which('soffice') or 'libreoffice'

def convert_word_to_pdf_com(file_path, output_path, orientation='auto'):
    """Converts DOCX to PDF using Word COM automation on Windows."""
    # Ensure absolute paths
    abs_input = os.path.abspath(file_path)
    abs_output = os.path.abspath(output_path)
    
    # We initialize COM and open Word
    import comtypes.client
    word = comtypes.client.CreateObject("Word.Application")
    word.Visible = False
    doc = None
    try:
        doc = word.Documents.Open(abs_input)
        if orientation.lower() == 'landscape':
            doc.PageSetup.Orientation = 1
        elif orientation.lower() == 'portrait':
            doc.PageSetup.Orientation = 0
        # 17 represents wdFormatPDF
        doc.SaveAs(abs_output, FileFormat=17)
    finally:
        if doc is not None:
            doc.Close()
        word.Quit()

def convert_excel_to_pdf_com(file_path, output_path, orientation='auto'):
    """Converts XLSX to PDF using Excel COM automation on Windows."""
    abs_input = os.path.abspath(file_path)
    abs_output = os.path.abspath(output_path)
    
    import win32com.client
    excel = win32com.client.Dispatch("Excel.Application")
    excel.Visible = False
    wb = None
    try:
        wb = excel.Workbooks.Open(abs_input)
        for sheet in wb.Sheets:
            if orientation.lower() == 'landscape':
                sheet.PageSetup.Orientation = 2
            elif orientation.lower() == 'portrait':
                sheet.PageSetup.Orientation = 1
        # Type=0 represents xlTypePDF
        wb.ExportAsFixedFormat(0, abs_output)
    finally:
        if wb is not None:
            wb.Close(False)
        excel.Quit()

def _word_to_pdf_reportlab_fallback(file_path, output_path, orientation='auto', layout_mode='fit'):
    """Legacy cross-platform DOCX→PDF renderer using python-docx + reportlab.
    Used only when LibreOffice is not available (e.g. stripped dev environments)."""
    import docx

    doc = docx.Document(file_path)
    if len(doc.paragraphs) > 1000:
        raise ValueError("Word document exceeds 1000 paragraph limit.")

    pagesize = letter
    if orientation.lower() == 'landscape':
        pagesize = landscape(letter)

    margin = 40
    if layout_mode.lower() == 'fit':
        margin = 30

    doc_template = SimpleDocTemplate(
        output_path,
        pagesize=pagesize,
        rightMargin=margin, leftMargin=margin, topMargin=margin, bottomMargin=margin
    )

    styles = getSampleStyleSheet()
    story = []

    heading1_style = ParagraphStyle(
        name='DocxHeading1', parent=styles['Heading1'],
        fontSize=18, leading=22, spaceAfter=8
    )
    heading2_style = ParagraphStyle(
        name='DocxHeading2', parent=styles['Heading2'],
        fontSize=14, leading=18, spaceAfter=6
    )
    normal_style = ParagraphStyle(
        name='DocxNormal', parent=styles['Normal'],
        fontSize=10, leading=14, spaceAfter=6
    )

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            story.append(Spacer(1, 10))
            continue
        style_name = paragraph.style.name.lower()
        if "heading 1" in style_name:
            story.append(Paragraph(text, heading1_style))
        elif "heading 2" in style_name:
            story.append(Paragraph(text, heading2_style))
        else:
            story.append(Paragraph(text, normal_style))

    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_data.append(row_data)
        if table_data:
            t = Table(table_data)
            t.setStyle(TableStyle([
                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
                ('PADDING', (0, 0), (-1, -1), 6),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
            ]))
            story.append(Spacer(1, 10))
            story.append(t)
            story.append(Spacer(1, 10))

    if not story:
        story.append(Paragraph("Empty Document", normal_style))

    doc_template.build(story)


def _apply_docx_orientation(file_path, orientation, work_dir):
    """Clone the DOCX into a temp location and patch its page orientation via
    python-docx so LibreOffice inherits the correct layout.  Returns the path
    to the (possibly modified) copy."""
    import docx
    from docx.oxml.ns import qn
    from lxml import etree

    dest = os.path.join(work_dir, os.path.basename(file_path))
    shutil.copy2(file_path, dest)

    if orientation.lower() not in ('landscape', 'portrait'):
        return dest  # 'auto' — leave document as-is

    doc = docx.Document(dest)
    new_orient = 'landscape' if orientation.lower() == 'landscape' else 'portrait'
    w_orient = 'landscape' if new_orient == 'landscape' else None  # None removes the attribute

    for section in doc.sections:
        pgSz = section._sectPr.find(qn('w:pgSz'))
        if pgSz is None:
            pgSz = etree.SubElement(section._sectPr, qn('w:pgSz'))
        if w_orient:
            pgSz.set(qn('w:orient'), w_orient)
            # Swap width/height so LibreOffice renders the correct page size
            w_val = pgSz.get(qn('w:w'))
            h_val = pgSz.get(qn('w:h'))
            if w_val and h_val:
                if int(w_val) < int(h_val):  # currently portrait — flip to landscape
                    pgSz.set(qn('w:w'), h_val)
                    pgSz.set(qn('w:h'), w_val)
        else:
            pgSz.attrib.pop(qn('w:orient'), None)
            w_val = pgSz.get(qn('w:w'))
            h_val = pgSz.get(qn('w:h'))
            if w_val and h_val:
                if int(w_val) > int(h_val):  # currently landscape — flip to portrait
                    pgSz.set(qn('w:w'), h_val)
                    pgSz.set(qn('w:h'), w_val)

    doc.save(dest)
    return dest


def word_to_pdf_portable(file_path, output_path, orientation='auto', layout_mode='fit'):
    """Convert DOCX → PDF using headless LibreOffice for professional layout fidelity.
    Falls back to the legacy reportlab renderer if LibreOffice is not installed."""

    if not LIBREOFFICE_AVAILABLE:
        # Local dev without LibreOffice — use legacy renderer
        _word_to_pdf_reportlab_fallback(file_path, output_path, orientation=orientation, layout_mode=layout_mode)
        return

    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    with tempfile.TemporaryDirectory() as work_dir:
        # Apply orientation settings to a temporary copy of the DOCX
        input_file = _apply_docx_orientation(file_path, orientation, work_dir)

        # LibreOffice writes <stem>.pdf into work_dir
        try:
            subprocess.run(
                [
                    _LO_BIN,
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', work_dir,
                    input_file,
                ],
                check=True,
                timeout=120,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
            )
        except subprocess.CalledProcessError as exc:
            stderr_msg = exc.stderr.decode(errors='replace') if exc.stderr else ''
            raise RuntimeError(f"LibreOffice conversion failed: {stderr_msg}") from exc

        # Locate the generated PDF (LibreOffice names it after the input stem)
        stem = os.path.splitext(os.path.basename(input_file))[0]
        lo_output = os.path.join(work_dir, stem + '.pdf')

        if not os.path.exists(lo_output):
            # Search for any PDF produced in work_dir as a safety net
            pdfs = [f for f in os.listdir(work_dir) if f.lower().endswith('.pdf')]
            if not pdfs:
                raise RuntimeError("LibreOffice did not produce a PDF output file.")
            lo_output = os.path.join(work_dir, pdfs[0])

        shutil.move(lo_output, output_path)

def excel_to_pdf_portable(file_path, output_path, orientation='auto', layout_mode='fit'):
    """Fallback cross-platform xlsx to pdf renderer using openpyxl and reportlab."""
    import openpyxl
    
    wb = openpyxl.load_workbook(file_path, data_only=True)
    try:
        if len(wb.sheetnames) > 20:
            raise ValueError("Excel file exceeds 20 sheet limit.")

        pagesize = landscape(letter)
        if orientation.lower() == 'portrait':
            pagesize = letter
            
        margin = 30
        if layout_mode.lower() == 'fit':
            margin = 15
            
        doc_template = SimpleDocTemplate(
            output_path, 
            pagesize=pagesize,
            rightMargin=margin, leftMargin=margin, topMargin=margin, bottomMargin=margin
        )
        
        styles = getSampleStyleSheet()
        story = []
        
        title_style = ParagraphStyle(
            name='SheetTitle', 
            parent=styles['Heading2'], 
            fontSize=12, 
            leading=16,
            spaceAfter=10
        )
        
        normal_style = ParagraphStyle(
            name='ExcelNormal',
            parent=styles['Normal'],
            fontSize=8,
            leading=10
        )
        
        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            sheet = wb[sheet_name]
            if sheet_idx > 0:
                story.append(PageBreak())
                
            story.append(Paragraph(f"Sheet: {sheet_name}", title_style))
            
            # Read grid data
            data = []
            for row in sheet.iter_rows(values_only=True):
                if not any(row is not None and str(row_cell).strip() for row_cell in row if row_cell is not None):
                    continue # Skip empty row
                row_data = [str(cell) if cell is not None else "" for cell in row]
                data.append(row_data)
                
            if not data:
                story.append(Paragraph("Empty Sheet", styles['Normal']))
                continue
                
            # Chop cols to prevent overflowing width
            max_cols = 10
            data_sliced = [r[:max_cols] for r in data]
            
            # Wrap strings in Paragraphs to auto-wrap inside table cells
            wrapped_data = []
            for r_idx, row_list in enumerate(data_sliced):
                wrapped_row = []
                for col_val in row_list:
                    wrapped_row.append(Paragraph(col_val, normal_style))
                wrapped_data.append(wrapped_row)
                
            t = Table(wrapped_data)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#0F62FE")),
                ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
                ('ALIGN', (0,0), (-1,-1), 'LEFT'),
                ('GRID', (0,0), (-1,-1), 0.5, colors.lightgrey),
                ('PADDING', (0,0), (-1,-1), 4),
                ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ]))
            story.append(t)
            
        doc_template.build(story)
    finally:
        wb.close()

def word_to_pdf(file_path, output_path, orientation='auto', layout_mode='fit'):
    if COM_AVAILABLE:
        try:
            convert_word_to_pdf_com(file_path, output_path, orientation=orientation)
            return output_path
        except Exception as com_err:
            # Fall back to python-docx + reportlab
            word_to_pdf_portable(file_path, output_path, orientation=orientation, layout_mode=layout_mode)
            return output_path
    else:
        word_to_pdf_portable(file_path, output_path, orientation=orientation, layout_mode=layout_mode)
        return output_path

def excel_to_pdf(file_path, output_path, orientation='auto', layout_mode='fit'):
    if COM_AVAILABLE:
        try:
            convert_excel_to_pdf_com(file_path, output_path, orientation=orientation)
            return output_path
        except Exception as com_err:
            excel_to_pdf_portable(file_path, output_path, orientation=orientation, layout_mode=layout_mode)
            return output_path
    else:
        excel_to_pdf_portable(file_path, output_path, orientation=orientation, layout_mode=layout_mode)
        return output_path

def pdf_to_word(file_path, output_path, mode='flowing', ocr=False):
    assert_pdf_limits(file_path)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    cv = Converter(file_path)
    try:
        # Convert all pages
        cv.convert(output_path, start=0, end=None)
    finally:
        cv.close()
    return output_path

def pdf_to_excel(file_path, output_path, mode='auto', single_sheet=False):
    import openpyxl
    import re
    
    reader = assert_pdf_limits(file_path)
    wb = openpyxl.Workbook()
    try:
        # Remove default sheet
        default_sheet = wb.active
        wb.remove(default_sheet)
        
        ws = None
        if single_sheet:
            ws = wb.create_sheet(title="Extracted Data")
            
        row_num = 1
        for idx, page in enumerate(reader.pages):
            if not single_sheet:
                ws = wb.create_sheet(title=f"Page {idx+1}")
                row_num = 1
                
            text = page.extract_text()
            if not text:
                continue
            
            for line in text.split('\n'):
                line = line.strip()
                if not line:
                    continue
                parts = re.split(r'\t| {2,}', line)
                for col_num, part in enumerate(parts):
                    ws.cell(row=row_num, column=col_num+1, value=part)
                row_num += 1
                
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        wb.save(output_path)
    finally:
        wb.close()
    return output_path

def pdf_to_ppt(file_path, output_path, slide_size='16:9', vector_mode=False):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    reader = assert_pdf_limits(file_path)
    prs = Presentation()
    
    # Set slide size
    if slide_size == '16:9':
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
    blank_layout = prs.slide_layouts[6] # Blank slide
    
    for idx, page in enumerate(reader.pages):
        slide = prs.slides.add_slide(blank_layout)
        text = page.extract_text()
        if not text:
            continue
        
        # Sizing box
        left = Inches(0.75)
        top = Inches(0.75)
        width = prs.slide_width - Inches(1.5)
        height = prs.slide_height - Inches(1.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, line in enumerate(text.split('\n')):
            line = line.strip()
            if not line:
                continue
            if i == 0:
                p = tf.paragraphs[0]
                p.text = line
                p.font.bold = True
                p.font.size = Pt(20)
                p.space_after = Pt(14)
            else:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(12)
                p.space_after = Pt(6)
                
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path

def convert_ppt_to_pdf_com(file_path, output_path):
    abs_input = os.path.abspath(file_path)
    abs_output = os.path.abspath(output_path)
    
    import win32com.client
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    pres = None
    try:
        pres = powerpoint.Presentations.Open(abs_input, WithWindow=False)
        # ppSaveAsPDF is 32
        pres.SaveAs(abs_output, FileFormat=32)
    finally:
        if pres is not None:
            pres.Close()
        powerpoint.Quit()

def _ppt_to_pdf_reportlab_fallback(file_path, output_path, orientation='auto', layout_mode='fit'):
    """Legacy cross-platform PPTX→PDF renderer using python-pptx + reportlab.
    Used only when LibreOffice is not available."""
    from pptx import Presentation

    prs = Presentation(file_path)
    if len(prs.slides) > 100:
        raise ValueError("PowerPoint presentation exceeds 100 slide limit.")

    pagesize = landscape(letter)
    if orientation.lower() == 'portrait':
        pagesize = letter

    margin = 40
    if layout_mode.lower() == 'fit':
        margin = 20

    doc_template = SimpleDocTemplate(
        output_path,
        pagesize=pagesize,
        rightMargin=margin, leftMargin=margin, topMargin=margin, bottomMargin=margin
    )

    styles = getSampleStyleSheet()
    story = []

    normal_style = ParagraphStyle(
        name='PptNormal', parent=styles['Normal'],
        fontSize=10, leading=14, spaceAfter=6
    )
    title_style = ParagraphStyle(
        name='PptTitle', parent=styles['Heading2'],
        fontSize=16, leading=20, spaceAfter=12
    )

    for idx, slide in enumerate(prs.slides):
        if idx > 0:
            story.append(PageBreak())
        story.append(Paragraph(f"Slide {idx+1}", title_style))
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        story.append(Paragraph(text, normal_style))

    if not story:
        story.append(Paragraph("Empty Presentation", normal_style))

    doc_template.build(story)


def ppt_to_pdf_portable(file_path, output_path, orientation='auto', layout_mode='fit'):
    """Convert PPTX → PDF using headless LibreOffice Impress for full slide fidelity.
    Falls back to the legacy reportlab renderer if LibreOffice is not installed."""

    if not LIBREOFFICE_AVAILABLE:
        _ppt_to_pdf_reportlab_fallback(file_path, output_path, orientation=orientation, layout_mode=layout_mode)
        return

    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    with tempfile.TemporaryDirectory() as work_dir:
        # Copy input so LibreOffice writes next to it in our controlled temp dir
        input_copy = os.path.join(work_dir, os.path.basename(file_path))
        shutil.copy2(file_path, input_copy)

        try:
            subprocess.run(
                [
                    _LO_BIN,
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', work_dir,
                    input_copy,
                ],
                check=True,
                timeout=180,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
            )
        except subprocess.CalledProcessError as exc:
            stderr_msg = exc.stderr.decode(errors='replace') if exc.stderr else ''
            raise RuntimeError(f"LibreOffice PPTX conversion failed: {stderr_msg}") from exc

        stem = os.path.splitext(os.path.basename(input_copy))[0]
        lo_output = os.path.join(work_dir, stem + '.pdf')

        if not os.path.exists(lo_output):
            pdfs = [f for f in os.listdir(work_dir) if f.lower().endswith('.pdf')]
            if not pdfs:
                raise RuntimeError("LibreOffice did not produce a PDF output file for PPTX.")
            lo_output = os.path.join(work_dir, pdfs[0])

        shutil.move(lo_output, output_path)

def ppt_to_pdf(file_path, output_path, orientation='auto', layout_mode='fit'):
    if COM_AVAILABLE:
        try:
            convert_ppt_to_pdf_com(file_path, output_path)
            return output_path
        except Exception:
            ppt_to_pdf_portable(file_path, output_path, orientation=orientation, layout_mode=layout_mode)
            return output_path
    else:
        ppt_to_pdf_portable(file_path, output_path, orientation=orientation, layout_mode=layout_mode)
        return output_path

from html.parser import HTMLParser

class SimpleHTMLParser(HTMLParser):
    def __init__(self, styles):
        super().__init__()
        self.styles = styles
        self.story = []
        self.current_tag = None
        self.current_text = ""
        
    def handle_starttag(self, tag, attrs):
        self.current_tag = tag
        self.current_text = ""
        
    def handle_endtag(self, tag):
        text = self.current_text.strip()
        if text:
            if tag == 'h1':
                self.story.append(Paragraph(text, self.styles['Heading1']))
                self.story.append(Spacer(1, 10))
            elif tag == 'h2':
                self.story.append(Paragraph(text, self.styles['Heading2']))
                self.story.append(Spacer(1, 8))
            elif tag == 'p':
                self.story.append(Paragraph(text, self.styles['Normal']))
                self.story.append(Spacer(1, 6))
            elif tag == 'li':
                self.story.append(Paragraph(f"• {text}", self.styles['Normal']))
                self.story.append(Spacer(1, 4))
        self.current_tag = None
        
    def handle_data(self, data):
        if self.current_tag:
            self.current_text += data
        else:
            text = data.strip()
            if text:
                self.story.append(Paragraph(text, self.styles['Normal']))
                self.story.append(Spacer(1, 6))

def html_to_pdf(file_path, output_path, orientation='auto', layout_mode='fit'):
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        html_content = f.read()
        
    pagesize = letter
    if orientation.lower() == 'landscape':
        pagesize = landscape(letter)
        
    margin = 40
    if layout_mode.lower() == 'fit':
        margin = 25
        
    doc_template = SimpleDocTemplate(
        output_path, 
        pagesize=pagesize,
        rightMargin=margin, leftMargin=margin, topMargin=margin, bottomMargin=margin
    )
    
    styles = getSampleStyleSheet()
    parser = SimpleHTMLParser(styles)
    parser.feed(html_content)
    
    if not parser.story:
        parser.story.append(Paragraph("Empty HTML Document", styles['Normal']))
        
    doc_template.build(parser.story)
    return output_path

def main():
    try:
        # Check command line arguments first
        if len(sys.argv) >= 3:
            action = sys.argv[1]
            if action in ["word2pdf", "excel2pdf", "ppt2pdf", "html2pdf"]:
                file_path = sys.argv[2]
                output = sys.argv[3]
                orientation = sys.argv[4] if len(sys.argv) > 4 else "auto"
                layout_mode = sys.argv[5] if len(sys.argv) > 5 else "fit"
                if action == "word2pdf":
                    result = word_to_pdf(file_path, output, orientation=orientation, layout_mode=layout_mode)
                elif action == "excel2pdf":
                    result = excel_to_pdf(file_path, output, orientation=orientation, layout_mode=layout_mode)
                elif action == "ppt2pdf":
                    result = ppt_to_pdf(file_path, output, orientation=orientation, layout_mode=layout_mode)
                else:
                    result = html_to_pdf(file_path, output, orientation=orientation, layout_mode=layout_mode)
                print(json.dumps({"success": True, "output": result}))
                return
            elif action == "pdf2word":
                file_path = sys.argv[2]
                output = sys.argv[3]
                mode = sys.argv[4] if len(sys.argv) > 4 else "flowing"
                ocr = (sys.argv[5].lower() == "true") if len(sys.argv) > 5 else False
                result = pdf_to_word(file_path, output, mode=mode, ocr=ocr)
                print(json.dumps({"success": True, "output": result}))
                return
            elif action == "pdf2excel":
                file_path = sys.argv[2]
                output = sys.argv[3]
                mode = sys.argv[4] if len(sys.argv) > 4 else "auto"
                single_sheet = (sys.argv[5].lower() == "true") if len(sys.argv) > 5 else False
                result = pdf_to_excel(file_path, output, mode=mode, single_sheet=single_sheet)
                print(json.dumps({"success": True, "output": result}))
                return
            elif action == "pdf2ppt":
                file_path = sys.argv[2]
                output = sys.argv[3]
                slide_size = sys.argv[4] if len(sys.argv) > 4 else "16:9"
                vector_mode = (sys.argv[5].lower() == "true") if len(sys.argv) > 5 else False
                result = pdf_to_ppt(file_path, output, slide_size=slide_size, vector_mode=vector_mode)
                print(json.dumps({"success": True, "output": result}))
                return

        # Read parameters from stdin
        input_data = sys.stdin.read()
        if not input_data:
            print(json.dumps({"success": False, "error": "No input arguments provided."}))
            return
            
        params = json.loads(input_data)
        action = params.get("action")
        settings = params.get("settings", {})
        
        if action == "word2pdf":
            file_path = params.get("file")
            output = params.get("output")
            orientation = settings.get("orientation", "auto")
            layout_mode = settings.get("layoutMode", "fit")
            result = word_to_pdf(file_path, output, orientation=orientation, layout_mode=layout_mode)
            print(json.dumps({"success": True, "output": result}))
            
        elif action == "excel2pdf":
            file_path = params.get("file")
            output = params.get("output")
            orientation = settings.get("orientation", "auto")
            layout_mode = settings.get("layoutMode", "fit")
            result = excel_to_pdf(file_path, output, orientation=orientation, layout_mode=layout_mode)
            print(json.dumps({"success": True, "output": result}))
            
        elif action == "pdf2word":
            file_path = params.get("file")
            output = params.get("output")
            mode = settings.get("mode", "flowing")
            ocr = settings.get("ocr", False)
            result = pdf_to_word(file_path, output, mode=mode, ocr=ocr)
            print(json.dumps({"success": True, "output": result}))

        elif action == "pdf2excel":
            file_path = params.get("file")
            output = params.get("output")
            mode = settings.get("mode", "auto")
            single_sheet = settings.get("singleSheet", False)
            result = pdf_to_excel(file_path, output, mode=mode, single_sheet=single_sheet)
            print(json.dumps({"success": True, "output": result}))

        elif action == "pdf2ppt":
            file_path = params.get("file")
            output = params.get("output")
            slide_size = settings.get("slideSize", "16:9")
            vector_mode = settings.get("vectorMode", False)
            result = pdf_to_ppt(file_path, output, slide_size=slide_size, vector_mode=vector_mode)
            print(json.dumps({"success": True, "output": result}))

        elif action == "ppt2pdf":
            file_path = params.get("file")
            output = params.get("output")
            orientation = settings.get("orientation", "auto")
            layout_mode = settings.get("layoutMode", "fit")
            result = ppt_to_pdf(file_path, output, orientation=orientation, layout_mode=layout_mode)
            print(json.dumps({"success": True, "output": result}))

        elif action == "html2pdf":
            file_path = params.get("file")
            output = params.get("output")
            orientation = settings.get("orientation", "auto")
            layout_mode = settings.get("layoutMode", "fit")
            result = html_to_pdf(file_path, output, orientation=orientation, layout_mode=layout_mode)
            print(json.dumps({"success": True, "output": result}))
            
        else:
            print(json.dumps({"success": False, "error": f"Unknown action: {action}"}))
            
    except Exception as e:
        print(json.dumps({"success": False, "error": str(e)}))

if __name__ == "__main__":
    main()
