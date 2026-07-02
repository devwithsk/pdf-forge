import os
import sys
import json
import subprocess
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

def create_mock_pdf(filename, text="Test PDF Content", pages=1):
    os.makedirs(os.path.dirname(filename) if os.path.dirname(filename) else ".", exist_ok=True)
    c = canvas.Canvas(filename, pagesize=letter)
    for p in range(pages):
        c.drawString(100, 500, f"{text} - Page {p+1}")
        c.showPage()
    c.save()
    print(f"Created mock PDF: {filename}")

def create_mock_docx(filename, text="Test Word Content"):
    try:
        import docx
        doc = docx.Document()
        doc.add_heading('Mock Document Header', level=1)
        doc.add_paragraph(text)
        doc.save(filename)
        print(f"Created mock DOCX: {filename}")
    except ImportError:
        print("python-docx not installed yet, skipping DOCX generation")

def create_mock_xlsx(filename, sheet_name="Sheet1"):
    try:
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = sheet_name
        ws['A1'] = "Header 1"
        ws['B1'] = "Header 2"
        ws['A2'] = "Row 1 Col 1"
        ws['B2'] = "Row 1 Col 2"
        wb.save(filename)
        print(f"Created mock XLSX: {filename}")
    except ImportError:
        print("openpyxl not installed yet, skipping XLSX generation")

def create_mock_image(filename):
    try:
        from PIL import Image, ImageDraw
        img = Image.new('RGB', (200, 100), color = (73, 109, 137))
        d = ImageDraw.Draw(img)
        d.text((10,10), "Test Image", fill=(255,255,0))
        img.save(filename)
        print(f"Created mock Image: {filename}")
    except ImportError:
        print("Pillow not installed yet, skipping Image generation")

def create_mock_pptx(filename, text="Test presentation content"):
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Mock Header Title"
        slide.placeholders[1].text = text
        prs.save(filename)
        print(f"Created mock PPTX: {filename}")
    except Exception as e:
        print(f"Failed to create PPTX: {e}")

def create_mock_html(filename):
    with open(filename, "w", encoding="utf-8") as f:
        f.write("<h1>Sample Heading 1</h1>\n<p>This is a paragraph of mock html content.</p>\n<li>List Item 1</li>\n<li>List Item 2</li>")
    print(f"Created mock HTML: {filename}")

def run_script_module(script_path, input_dict):
    p = subprocess.Popen(
        [sys.executable, script_path],
        stdin=subprocess.PIPE,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True
    )
    stdout, stderr = p.communicate(input=json.dumps(input_dict))
    
    if p.returncode != 0:
        return {"success": False, "error": f"Process exited with {p.returncode}. Stderr: {stderr}"}
    
    try:
        return json.loads(stdout.strip())
    except Exception as e:
        return {"success": False, "error": f"Failed to parse stdout as JSON. Stdout: {stdout}. Error: {str(e)}"}

def main():
    print("Starting Python Engine Integration Tests...")
    
    # 1. Setup folders
    test_dir = os.path.dirname(os.path.abspath(__file__))
    output_dir = os.path.join(test_dir, "outputs")
    os.makedirs(output_dir, exist_ok=True)
    
    # Module paths
    modules_dir = os.path.join(os.path.dirname(test_dir), "modules")
    basic_script = os.path.join(modules_dir, "basic_manipulation.py")
    security_script = os.path.join(modules_dir, "security.py")
    image_script = os.path.join(modules_dir, "image_convert.py")
    doc_script = os.path.join(modules_dir, "doc_convert.py")
    
    # Files
    file_pdf1 = os.path.join(test_dir, "test1.pdf")
    file_pdf2 = os.path.join(test_dir, "test2.pdf")
    file_docx = os.path.join(test_dir, "test.docx")
    file_xlsx = os.path.join(test_dir, "test.xlsx")
    file_png = os.path.join(test_dir, "test.png")
    file_pptx = os.path.join(test_dir, "test.pptx")
    file_html = os.path.join(test_dir, "test.html")
    
    create_mock_pdf(file_pdf1, "First Document", pages=2)
    create_mock_pdf(file_pdf2, "Second Document", pages=1)
    create_mock_docx(file_docx, "Hello World from Docx file!")
    create_mock_xlsx(file_xlsx)
    create_mock_image(file_png)
    create_mock_pptx(file_pptx, "Hello slides content!")
    create_mock_html(file_html)
    
    # Test Merge
    print("\n--- Testing PDF Merge ---")
    out_merge = os.path.join(output_dir, "merged.pdf")
    res = run_script_module(basic_script, {
        "action": "merge",
        "files": [file_pdf1, file_pdf2],
        "output": out_merge
    })
    print("Merge Result:", res)
    assert res.get("success"), f"Merge failed: {res.get('error')}"
    assert os.path.exists(out_merge), "Merged PDF file does not exist!"
    
    # Test Split
    print("\n--- Testing PDF Split ---")
    res = run_script_module(basic_script, {
        "action": "split",
        "file": out_merge,
        "output_dir": output_dir,
        "split_mode": "all"
    })
    print("Split All Result:", res)
    assert res.get("success"), f"Split failed: {res.get('error')}"
    
    # Test Rotate
    print("\n--- Testing PDF Rotate ---")
    out_rotate = os.path.join(output_dir, "rotated.pdf")
    res = run_script_module(basic_script, {
        "action": "rotate",
        "file": file_pdf1,
        "output": out_rotate,
        "degrees": 90
    })
    print("Rotate Result:", res)
    assert res.get("success"), f"Rotate failed: {res.get('error')}"
    
    # Test Protect
    print("\n--- Testing PDF Protect ---")
    out_protect = os.path.join(output_dir, "protected.pdf")
    res = run_script_module(security_script, {
        "action": "protect",
        "file": file_pdf1,
        "output": out_protect,
        "password": "testpassword"
    })
    print("Protect Result:", res)
    assert res.get("success"), f"Protect failed: {res.get('error')}"
    
    # Test Unlock
    print("\n--- Testing PDF Unlock ---")
    out_unlock = os.path.join(output_dir, "unlocked.pdf")
    res = run_script_module(security_script, {
        "action": "unlock",
        "file": out_protect,
        "output": out_unlock,
        "password": "testpassword"
    })
    print("Unlock Result:", res)
    assert res.get("success"), f"Unlock failed: {res.get('error')}"
    
    # Test Watermark
    print("\n--- Testing PDF Watermark ---")
    out_watermark = os.path.join(output_dir, "watermarked.pdf")
    res = run_script_module(security_script, {
        "action": "watermark",
        "file": file_pdf1,
        "output": out_watermark,
        "text": "DRAFT",
        "font_size": 50,
        "opacity": 0.25,
        "color": "#FF0000"
    })
    print("Watermark Result:", res)
    assert res.get("success"), f"Watermark failed: {res.get('error')}"
    
    # Test JPG to PDF
    print("\n--- Testing JPG to PDF ---")
    out_jpg2pdf = os.path.join(output_dir, "jpg2pdf.pdf")
    res = run_script_module(image_script, {
        "action": "jpg2pdf",
        "images": [file_png],
        "output": out_jpg2pdf
    })
    print("JPG to PDF Result:", res)
    assert res.get("success"), f"JPG to PDF failed: {res.get('error')}"
    
    # Test DOCX to PDF
    print("\n--- Testing DOCX to PDF ---")
    out_docx2pdf = os.path.join(output_dir, "docx2pdf.pdf")
    res = run_script_module(doc_script, {
        "action": "word2pdf",
        "file": file_docx,
        "output": out_docx2pdf
    })
    print("DOCX to PDF Result:", res)
    assert res.get("success"), f"DOCX to PDF failed: {res.get('error')}"
    
    # Test XLSX to PDF
    print("\n--- Testing XLSX to PDF ---")
    out_xlsx2pdf = os.path.join(output_dir, "xlsx2pdf.pdf")
    res = run_script_module(doc_script, {
        "action": "excel2pdf",
        "file": file_xlsx,
        "output": out_xlsx2pdf
    })
    print("XLSX to PDF Result:", res)
    assert res.get("success"), f"XLSX to PDF failed: {res.get('error')}"
    
    # Test PDF to DOCX
    print("\n--- Testing PDF to DOCX ---")
    out_pdf2docx = os.path.join(output_dir, "pdf2docx.docx")
    res = run_script_module(doc_script, {
        "action": "pdf2word",
        "file": file_pdf2,
        "output": out_pdf2docx
    })
    print("PDF to DOCX Result:", res)
    assert res.get("success"), f"PDF to DOCX failed: {res.get('error')}"

    # Test PDF to XLSX
    print("\n--- Testing PDF to XLSX ---")
    out_pdf2xlsx = os.path.join(output_dir, "pdf2excel.xlsx")
    res = run_script_module(doc_script, {
        "action": "pdf2excel",
        "file": file_pdf1,
        "output": out_pdf2xlsx
    })
    print("PDF to XLSX Result:", res)
    assert res.get("success"), f"PDF to XLSX failed: {res.get('error')}"

    # Test PDF to PPTX
    print("\n--- Testing PDF to PPTX ---")
    out_pdf2pptx = os.path.join(output_dir, "pdf2ppt.pptx")
    res = run_script_module(doc_script, {
        "action": "pdf2ppt",
        "file": file_pdf1,
        "output": out_pdf2pptx
    })
    print("PDF to PPTX Result:", res)
    assert res.get("success"), f"PDF to PPTX failed: {res.get('error')}"

    # Test PPTX to PDF
    print("\n--- Testing PPTX to PDF ---")
    out_ppt2pdf = os.path.join(output_dir, "ppt2pdf.pdf")
    res = run_script_module(doc_script, {
        "action": "ppt2pdf",
        "file": file_pptx,
        "output": out_ppt2pdf
    })
    print("PPTX to PDF Result:", res)
    assert res.get("success"), f"PPTX to PDF failed: {res.get('error')}"

    # Test HTML to PDF
    print("\n--- Testing HTML to PDF ---")
    out_html2pdf = os.path.join(output_dir, "html2pdf.pdf")
    res = run_script_module(doc_script, {
        "action": "html2pdf",
        "file": file_html,
        "output": out_html2pdf
    })
    print("HTML to PDF Result:", res)
    assert res.get("success"), f"HTML to PDF failed: {res.get('error')}"
    
    print("\n===============================")
    print("ALL INTEGRATION TESTS PASSED!!!")
    print("===============================")

if __name__ == "__main__":
    main()
